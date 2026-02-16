#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate-pptx.py - 将论文精读报告转换为专业PPT幻灯片（增强版）

功能：
- 解析Markdown格式的精读报告
- 从PDF提取图片并插入PPT
- 渲染LaTeX公式为图片并插入PPT
- 生成25-30页专业PPT幻灯片
- 白色背景主题，蓝白橙配色
- 结构化内容分页，包含图表和公式

使用方法：
python3 generate-pptx.py input_report.md output_slides.pptx [--pdf paper.pdf]

依赖：python-pptx, PyMuPDF/fitz, matplotlib, Pillow

作者：太森的AI助手二丫
版本：v3.0 - 图片公式增强版
"""

import sys
import re
import os
import argparse
from datetime import datetime
from io import BytesIO
import fitz  # PyMuPDF
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class PPTGenerator:
    def __init__(self):
        # 配色方案
        self.colors = {
            'background': RGBColor(255, 255, 255),      # 白色 #FFFFFF
            'title': RGBColor(27, 58, 92),              # 深蓝 #1B3A5C
            'text': RGBColor(51, 51, 51),               # 深灰 #333333
            'accent': RGBColor(230, 126, 34),           # 橙色 #E67E22
            'secondary': RGBColor(52, 152, 219),        # 浅蓝 #3498DB
            'light_gray': RGBColor(248, 249, 250)       # 浅灰 #F8F9FA
        }
        
        # 字体大小
        self.font_sizes = {
            'title': Pt(28),
            'subtitle': Pt(24),
            'heading': Pt(20),
            'text': Pt(18),
            'caption': Pt(14),
            'small': Pt(12)
        }
        
        # 创建演示文稿
        self.prs = Presentation()
        self._setup_slide_master()
        
        # 存储PDF图片
        self.pdf_figures = []
        
        # 配置matplotlib中文字体
        plt.rcParams['font.family'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
        plt.rcParams['axes.unicode_minus'] = False

    def _setup_slide_master(self):
        """设置幻灯片母版样式"""
        # 设置幻灯片尺寸为16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def extract_figures_from_pdf(self, pdf_path):
        """从PDF提取图片"""
        if not os.path.exists(pdf_path):
            print(f"警告：PDF文件 {pdf_path} 不存在")
            return []
        
        print(f"正在从PDF提取图片: {pdf_path}")
        figures = []
        
        try:
            doc = fitz.open(pdf_path)
            for page_num in range(doc.page_count):
                page = doc[page_num]
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        # 转换CMYK到RGB
                        if pix.n >= 5:  # CMYK
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        
                        # 过滤太小的图片
                        if pix.width > 100 and pix.height > 100:
                            img_bytes = pix.tobytes("png")
                            figures.append({
                                'page': page_num + 1,  # 从1开始
                                'data': img_bytes,
                                'index': img_index,
                                'width': pix.width,
                                'height': pix.height
                            })
                            print(f"  提取图片: 第{page_num + 1}页, {pix.width}x{pix.height}")
                        
                        pix = None  # 释放内存
                    except Exception as e:
                        print(f"  跳过图片 (页{page_num + 1}, 索引{img_index}): {e}")
            
            doc.close()
            print(f"✅ 成功提取 {len(figures)} 张图片")
            
        except Exception as e:
            print(f"❌ PDF图片提取失败: {e}")
        
        return figures

    def render_formula(self, latex_str, fontsize=16):
        """渲染LaTeX公式为图片"""
        try:
            fig = plt.figure(figsize=(8, 1.5))
            fig.patch.set_facecolor('white')
            
            # 清理LaTeX字符串
            latex_clean = latex_str.strip()
            if latex_clean.startswith('$') and latex_clean.endswith('$'):
                latex_clean = latex_clean[1:-1]
            
            fig.text(0.5, 0.5, f'${latex_clean}$', fontsize=fontsize, 
                    ha='center', va='center', color='black')
            
            buf = BytesIO()
            fig.savefig(buf, format='png', bbox_inches='tight', 
                       dpi=150, facecolor='white', edgecolor='none')
            buf.seek(0)
            plt.close(fig)
            
            return buf.getvalue()
        except Exception as e:
            print(f"⚠️  LaTeX渲染失败: {latex_str[:50]}... - {e}")
            return None

    def add_image_to_slide(self, slide, image_bytes, left=None, top=None, width=None):
        """向幻灯片添加图片"""
        try:
            # 默认位置和大小
            if left is None:
                left = Inches(8)  # 右侧位置
            if top is None:
                top = Inches(2)
            if width is None:
                width = Inches(4)
            
            image_stream = BytesIO(image_bytes)
            slide.shapes.add_picture(image_stream, left, top, width)
            return True
        except Exception as e:
            print(f"⚠️  图片插入失败: {e}")
            return False

    def extract_latex_formulas(self, text):
        """从文本中提取LaTeX公式"""
        formulas = []
        
        # 匹配 $...$ 和 $$...$$ 
        patterns = [
            r'\$\$([^$]+)\$\$',  # 显示公式 $$...$$
            r'\$([^$]+)\$'       # 行内公式 $...$
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text)
            for match in matches:
                latex_content = match.group(1).strip()
                if len(latex_content) > 2:  # 过滤太短的
                    formulas.append({
                        'latex': latex_content,
                        'full_match': match.group(0),
                        'start': match.start(),
                        'end': match.end()
                    })
        
        return formulas

    def parse_markdown_report(self, content):
        """解析Markdown格式的精读报告"""
        sections = {}
        current_section = None
        current_content = []
        
        lines = content.split('\n')
        for line in lines:
            # 检测一级标题
            if line.startswith('# '):
                if current_section:
                    sections[current_section] = '\n'.join(current_content)
                current_section = 'title'
                current_content = [line[2:].strip()]
            # 检测二级标题
            elif line.startswith('## '):
                if current_section:
                    sections[current_section] = '\n'.join(current_content)
                section_name = line[3:].strip()
                current_section = self._normalize_section_name(section_name)
                current_content = []
            else:
                current_content.append(line)
        
        # 保存最后一个section
        if current_section:
            sections[current_section] = '\n'.join(current_content)
        
        return sections

    def _normalize_section_name(self, section):
        """标准化section名称"""
        section_map = {
            '核心贡献': 'contribution',
            '问题背景与动机': 'background',
            '方法详解': 'method',
            '实验结果分析': 'experiment',
            '五专家会诊': 'experts',
            '综合评分': 'score',
            '核心学习要点': 'takeaways',
            '医疗机器人迁移路径': 'medical',
            '推荐行动': 'action'
        }
        
        for key, value in section_map.items():
            if key in section:
                return value
        return 'other'

    def clean_markdown_format(self, text):
        """清理Markdown格式标记"""
        # 简化处理，避免复杂的格式标记系统
        text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # 移除粗体标记
        text = re.sub(r'\*([^*]+)\*', r'\1', text)      # 移除斜体标记
        text = re.sub(r'`([^`]+)`', r'\1', text)        # 移除代码标记
        text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'\1', text)  # 移除链接但保留文字
        text = re.sub(r'^\s*[-*+]\s*', '• ', text, flags=re.MULTILINE)  # 转换列表项
        text = re.sub(r'^\s*(\d+)\.\s*', r'\1. ', text, flags=re.MULTILINE)  # 保留数字列表
        text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)  # 移除标题标记
        text = re.sub(r'\n\s*\n+', '\n\n', text)  # 规范化空行
        
        return text.strip()

    def extract_paper_info(self, sections):
        """从sections中提取论文基本信息"""
        info = {
            'title_cn': '',
            'title_en': '',
            'authors': '',
            'institutions': '',
            'arxiv_link': '',
            'date': '',
            'category': '',
            'contribution': ''
        }
        
        title_section = sections.get('title', '')
        if title_section:
            lines = title_section.split('\n')
            for line in lines:
                if 'arxiv.org/abs/' in line.lower():
                    match = re.search(r'https?://arxiv\.org/abs/([^)\s]+)', line)
                    if match:
                        info['arxiv_link'] = f"https://arxiv.org/abs/{match.group(1)}"
                elif '英文标题' in line:
                    info['title_en'] = line.split(':', 1)[1].strip() if ':' in line else ''
                elif '作者团队' in line:
                    info['authors'] = line.split(':', 1)[1].strip() if ':' in line else ''
                elif '实验室' in line or '机构' in line:
                    info['institutions'] = line.split(':', 1)[1].strip() if ':' in line else ''
                elif '发表日期' in line:
                    info['date'] = line.split(':', 1)[1].strip() if ':' in line else ''
                elif '分类' in line:
                    info['category'] = line.split(':', 1)[1].strip() if ':' in line else ''
                else:
                    if not info['title_cn']:
                        info['title_cn'] = line.strip()
        
        # 获取核心贡献
        contribution = sections.get('contribution', '')
        if contribution:
            info['contribution'] = self.clean_markdown_format(contribution)[:200] + '...' if len(contribution) > 200 else self.clean_markdown_format(contribution)
        
        return info

    def add_title_slide(self, paper_info):
        """添加封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = paper_info['title_cn'] or '论文深度解读'
        
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = self.font_sizes['title']
        title_run.font.color.rgb = self.colors['title']
        title_run.font.bold = True
        
        # 英文标题
        if paper_info['title_en']:
            en_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
            en_title_frame = en_title_box.text_frame
            en_title_frame.text = paper_info['title_en']
            
            en_para = en_title_frame.paragraphs[0]
            en_para.alignment = PP_ALIGN.CENTER
            en_run = en_para.runs[0]
            en_run.font.size = self.font_sizes['subtitle']
            en_run.font.color.rgb = self.colors['text']
        
        # 作者和机构信息
        info_y = 3.5
        if paper_info['authors']:
            author_box = slide.shapes.add_textbox(Inches(0.5), Inches(info_y), Inches(12.33), Inches(0.5))
            author_frame = author_box.text_frame
            author_frame.text = f"作者：{paper_info['authors']}"
            author_para = author_frame.paragraphs[0]
            author_para.alignment = PP_ALIGN.CENTER
            author_run = author_para.runs[0]
            author_run.font.size = self.font_sizes['text']
            author_run.font.color.rgb = self.colors['text']
            info_y += 0.5
        
        if paper_info['institutions']:
            inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(info_y), Inches(12.33), Inches(0.5))
            inst_frame = inst_box.text_frame
            inst_frame.text = f"机构：{paper_info['institutions']}"
            inst_para = inst_frame.paragraphs[0]
            inst_para.alignment = PP_ALIGN.CENTER
            inst_run = inst_para.runs[0]
            inst_run.font.size = self.font_sizes['text']
            inst_run.font.color.rgb = self.colors['text']
            info_y += 0.5
        
        # arXiv链接
        if paper_info['arxiv_link']:
            link_box = slide.shapes.add_textbox(Inches(0.5), Inches(info_y), Inches(12.33), Inches(0.5))
            link_frame = link_box.text_frame
            link_frame.text = paper_info['arxiv_link']
            link_para = link_frame.paragraphs[0]
            link_para.alignment = PP_ALIGN.CENTER
            link_run = link_para.runs[0]
            link_run.font.size = self.font_sizes['text']
            link_run.font.color.rgb = self.colors['secondary']
            info_y += 0.7
        
        # 核心贡献
        if paper_info['contribution']:
            contrib_box = slide.shapes.add_textbox(Inches(1), Inches(info_y), Inches(11.33), Inches(1.5))
            contrib_frame = contrib_box.text_frame
            contrib_frame.text = f"核心贡献：{paper_info['contribution']}"
            contrib_para = contrib_frame.paragraphs[0]
            contrib_para.alignment = PP_ALIGN.CENTER
            contrib_run = contrib_para.runs[0]
            contrib_run.font.size = self.font_sizes['text']
            contrib_run.font.color.rgb = self.colors['accent']
        
        # 底部标识
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"AI论文深度解读 | {datetime.now().strftime('%Y-%m-%d')}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_run = footer_para.runs[0]
        footer_run.font.size = self.font_sizes['caption']
        footer_run.font.color.rgb = self.colors['text']

    def add_content_slide(self, title, content, slide_type='normal', pdf_figures=None, figure_index_start=0):
        """添加内容页，支持公式渲染和图片插入"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_run = title_para.runs[0]
        title_run.font.size = self.font_sizes['subtitle']
        title_run.font.color.rgb = self.colors['title']
        title_run.font.bold = True
        
        # 添加装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.5), Inches(1.0), 
            Inches(12.33), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.color.rgb = self.colors['accent']
        
        # 检测并渲染LaTeX公式
        formulas = self.extract_latex_formulas(content)
        
        # 有图片或公式时，调整布局
        has_visual_content = bool(formulas) or bool(pdf_figures)
        if has_visual_content:
            text_width = Inches(7)  # 左侧文字区域
            visual_x = Inches(8)    # 右侧图片/公式区域
        else:
            text_width = Inches(12.33)
        
        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), text_width, Inches(5.7))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # 处理文本内容（移除LaTeX公式）
        text_content = content
        for formula in formulas:
            text_content = text_content.replace(formula['full_match'], '[公式]')
        
        # 处理不同类型的内容
        if slide_type == 'bullet_list':
            self._add_bullet_content(content_frame, text_content)
        elif slide_type == 'two_column':
            self._add_two_column_content(slide, text_content, 1.3)
        else:
            self._add_normal_content(content_frame, text_content)
        
        # 添加公式图片
        visual_y = Inches(1.5)
        if formulas:
            for i, formula in enumerate(formulas[:3]):  # 最多3个公式
                formula_img = self.render_formula(formula['latex'])
                if formula_img:
                    self.add_image_to_slide(slide, formula_img, visual_x, visual_y, Inches(3.5))
                    visual_y += Inches(1.5)
        
        # 添加PDF图片
        if pdf_figures and figure_index_start < len(pdf_figures):
            figure = pdf_figures[figure_index_start]
            self.add_image_to_slide(slide, figure['data'], visual_x, visual_y, Inches(4))

    def _add_normal_content(self, text_frame, content):
        """添加普通文本内容（修复API问题）"""
        text_frame.clear()  # 清空默认段落
        
        paragraphs = content.split('\n\n')
        for i, para in enumerate(paragraphs):
            if para.strip():
                if i == 0:
                    p = text_frame.paragraphs[0]  # 使用默认段落
                else:
                    p = text_frame.add_paragraph()  # 正确的API
                
                p.text = para.strip()
                p.font.size = self.font_sizes['text']
                p.font.color.rgb = self.colors['text']

    def _add_bullet_content(self, text_frame, content):
        """添加项目符号内容"""
        text_frame.clear()
        
        lines = content.split('\n')
        for i, line in enumerate(lines):
            if line.strip():
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # 设置项目符号
                if line.strip().startswith('•') or line.strip().startswith('-'):
                    p.level = 0
                    text = line.strip()[1:].strip()
                elif line.strip().startswith('  •') or line.strip().startswith('  -'):
                    p.level = 1
                    text = line.strip()[3:].strip()
                else:
                    p.level = 0
                    text = line.strip()
                
                p.text = text
                p.font.size = self.font_sizes['text']
                p.font.color.rgb = self.colors['text']

    def _add_two_column_content(self, slide, content, start_y):
        """添加两栏内容"""
        parts = content.split('\n---\n')  # 用---分隔两栏
        
        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), Inches(6), Inches(5.7))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        if len(parts) > 0:
            self._add_normal_content(left_frame, parts[0])
        
        # 右栏
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(start_y), Inches(6), Inches(5.7))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        if len(parts) > 1:
            self._add_normal_content(right_frame, parts[1])

    def add_background_slides(self, sections, pdf_figures=None):
        """添加问题背景页面"""
        background_content = sections.get('background', '')
        if not background_content:
            return
        
        content = self.clean_markdown_format(background_content)
        
        # 分割内容为多个页面
        paragraphs = content.split('\n\n')
        current_page_content = []
        page_count = 0
        figure_index = 0
        
        for para in paragraphs:
            if para.strip():
                current_page_content.append(para.strip())
                
                # 每页大约2-3个段落
                if len(current_page_content) >= 2 and len('\n\n'.join(current_page_content)) > 800:
                    page_count += 1
                    if page_count == 1:
                        title = "问题背景与动机"
                    else:
                        title = f"问题背景与动机 ({page_count})"
                    
                    page_content = '\n\n'.join(current_page_content)
                    self.add_content_slide(title, page_content, 'normal', pdf_figures, figure_index)
                    current_page_content = []
                    if pdf_figures and figure_index < len(pdf_figures):
                        figure_index += 1
        
        # 处理剩余内容
        if current_page_content:
            page_count += 1
            if page_count == 1:
                title = "问题背景与动机"
            else:
                title = f"问题背景与动机 ({page_count})"
            
            page_content = '\n\n'.join(current_page_content)
            self.add_content_slide(title, page_content, 'normal', pdf_figures, figure_index)

    def add_method_slides(self, sections, pdf_figures=None):
        """添加核心方法页面"""
        method_content = sections.get('method', '')
        if not method_content:
            return
        
        content = self.clean_markdown_format(method_content)
        
        # 检测子sections（通过###标题）
        subsections = re.split(r'\n### ([^\n]+)\n', content)
        figure_index = 1  # 从第二张图开始用于方法页面
        
        if len(subsections) > 1:
            # 有子sections，每个子section一页
            if subsections[0].strip():
                self.add_content_slide("核心方法概览", subsections[0], 'normal', pdf_figures, figure_index)
                if pdf_figures and figure_index < len(pdf_figures):
                    figure_index += 1
            
            for i in range(1, len(subsections), 2):
                if i + 1 < len(subsections):
                    subsection_title = f"核心方法：{subsections[i]}"
                    subsection_content = subsections[i + 1]
                    self.add_content_slide(subsection_title, subsection_content, 'normal', pdf_figures, figure_index)
                    if pdf_figures and figure_index < len(pdf_figures):
                        figure_index += 1
        else:
            # 没有子sections，按段落分页
            paragraphs = content.split('\n\n')
            current_page_content = []
            page_count = 0
            
            for para in paragraphs:
                if para.strip():
                    current_page_content.append(para.strip())
                    
                    # 每页控制内容长度
                    if len('\n\n'.join(current_page_content)) > 1000:
                        page_count += 1
                        if page_count == 1:
                            title = "核心方法"
                        else:
                            title = f"核心方法 ({page_count})"
                        
                        page_content = '\n\n'.join(current_page_content)
                        self.add_content_slide(title, page_content, 'normal', pdf_figures, figure_index)
                        current_page_content = []
                        if pdf_figures and figure_index < len(pdf_figures):
                            figure_index += 1
            
            # 处理剩余内容
            if current_page_content:
                page_count += 1
                if page_count == 1:
                    title = "核心方法"
                else:
                    title = f"核心方法 ({page_count})"
                
                page_content = '\n\n'.join(current_page_content)
                self.add_content_slide(title, page_content, 'normal', pdf_figures, figure_index)

    def add_experiment_slides(self, sections, pdf_figures=None):
        """添加实验结果页面"""
        experiment_content = sections.get('experiment', '')
        if not experiment_content:
            return
        
        content = self.clean_markdown_format(experiment_content)
        figure_index = min(len(pdf_figures) // 2, 5) if pdf_figures else 0  # 实验用中后部分图片
        
        # 尝试识别表格数据
        if '|' in content:
            # 包含表格，特殊处理
            parts = content.split('\n\n')
            table_parts = []
            text_parts = []
            
            for part in parts:
                if '|' in part and part.count('|') > 3:
                    table_parts.append(part)
                else:
                    text_parts.append(part)
            
            # 文字内容页
            if text_parts:
                text_content = '\n\n'.join(text_parts)
                self.add_content_slide("实验设置与分析", text_content, 'normal', pdf_figures, figure_index)
                if pdf_figures and figure_index < len(pdf_figures):
                    figure_index += 1
            
            # 表格数据页
            if table_parts:
                table_content = '\n\n'.join(table_parts)
                self.add_content_slide("实验结果数据", table_content, 'two_column', pdf_figures, figure_index)
        else:
            # 普通文本，按段落分页
            paragraphs = content.split('\n\n')
            current_page_content = []
            page_count = 0
            
            for para in paragraphs:
                if para.strip():
                    current_page_content.append(para.strip())
                    
                    if len('\n\n'.join(current_page_content)) > 900:
                        page_count += 1
                        if page_count == 1:
                            title = "实验结果分析"
                        else:
                            title = f"实验结果分析 ({page_count})"
                        
                        page_content = '\n\n'.join(current_page_content)
                        self.add_content_slide(title, page_content, 'normal', pdf_figures, figure_index)
                        current_page_content = []
                        if pdf_figures and figure_index < len(pdf_figures):
                            figure_index += 1
            
            # 处理剩余内容
            if current_page_content:
                page_count += 1
                if page_count == 1:
                    title = "实验结果分析"
                else:
                    title = f"实验结果分析 ({page_count})"
                
                page_content = '\n\n'.join(current_page_content)
                self.add_content_slide(title, page_content, 'normal', pdf_figures, figure_index)

    def add_expert_slides(self, sections):
        """添加五专家会诊页面"""
        experts_content = sections.get('experts', '')
        if not experts_content:
            return
        
        content = self.clean_markdown_format(experts_content)
        
        # 识别五个专家的评价
        expert_sections = re.split(r'### ([^#\n]+专家评分[^#\n]*)', content)
        
        if len(expert_sections) > 1:
            # 每个专家一页
            for i in range(1, len(expert_sections), 2):
                if i + 1 < len(expert_sections):
                    expert_title = expert_sections[i].strip()
                    expert_content = expert_sections[i + 1].strip()
                    
                    # 清理标题格式
                    clean_title = re.sub(r'评分[：:]\s*\d+/10', '', expert_title)
                    
                    self.add_content_slide(clean_title, expert_content)
        else:
            # 没有明确分段，直接作为一页
            self.add_content_slide("专家会诊评价", content)

    def add_medical_slides(self, sections):
        """添加医疗机器人迁移页面"""
        medical_content = sections.get('medical', '')
        if not medical_content:
            return
        
        content = self.clean_markdown_format(medical_content)
        
        # 分为应用场景和迁移方案
        if len(content) > 1200:
            parts = content.split('\n\n')
            mid = len(parts) // 2
            
            part1 = '\n\n'.join(parts[:mid])
            part2 = '\n\n'.join(parts[mid:])
            
            self.add_content_slide("医疗机器人应用场景", part1)
            self.add_content_slide("技术迁移方案与挑战", part2)
        else:
            self.add_content_slide("医疗机器人迁移路径", content)

    def add_takeaways_slide(self, sections):
        """添加关键要点页面"""
        takeaways_content = sections.get('takeaways', '')
        if not takeaways_content:
            return
        
        content = self.clean_markdown_format(takeaways_content)
        
        # 转换为项目符号格式
        lines = content.split('\n')
        bullet_content = []
        
        for line in lines:
            if line.strip() and not line.startswith('#'):
                if not line.startswith('•') and not line.startswith('-'):
                    bullet_content.append(f"• {line.strip()}")
                else:
                    bullet_content.append(line.strip())
        
        formatted_content = '\n'.join(bullet_content)
        self.add_content_slide("核心学习要点", formatted_content, 'bullet_list')

    def add_action_slide(self, sections):
        """添加行动建议页面"""
        action_content = sections.get('action', '')
        if not action_content:
            # 生成通用行动建议
            action_content = """• 深入阅读论文原文，特别关注技术细节部分

• 查看作者团队的其他相关工作

• 尝试复现核心实验结果

• 思考如何将技术应用到自己的研究领域

• 关注后续相关工作的发展"""
        
        content = self.clean_markdown_format(action_content)
        
        # 确保是项目符号格式
        if '•' not in content:
            lines = content.split('\n')
            bullet_lines = []
            for line in lines:
                if line.strip():
                    bullet_lines.append(f"• {line.strip()}")
            content = '\n'.join(bullet_lines)
        
        self.add_content_slide("推荐行动清单", content, 'bullet_list')

    def add_conclusion_slide(self):
        """添加结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['light_gray']
        
        # 感谢文字
        thanks_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(2))
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = "感谢观看\n\n欢迎讨论与交流"
        
        for para in thanks_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = self.font_sizes['title']
            run.font.color.rgb = self.colors['title']
            run.font.bold = True
        
        # 底部信息
        footer_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.33), Inches(1))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"AI论文深度解读 | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_run = footer_para.runs[0] if footer_para.runs else footer_para.add_run()
        footer_run.font.size = self.font_sizes['text']
        footer_run.font.color.rgb = self.colors['text']

    def generate_pptx(self, sections, pdf_path=None):
        """生成完整的PPT"""
        # 提取PDF图片
        pdf_figures = []
        if pdf_path:
            pdf_figures = self.extract_figures_from_pdf(pdf_path)
        
        # 提取论文信息
        paper_info = self.extract_paper_info(sections)
        
        # 1. 封面页
        self.add_title_slide(paper_info)
        
        # 2. 问题背景（2-3页）
        self.add_background_slides(sections, pdf_figures)
        
        # 3. 核心方法（10-15页，重点含图含公式）
        self.add_method_slides(sections, pdf_figures)
        
        # 4. 实验结果（3-5页）
        self.add_experiment_slides(sections, pdf_figures)
        
        # 5. 专家会诊（5页）
        self.add_expert_slides(sections)
        
        # 6. 医疗机器人迁移（2-3页）
        self.add_medical_slides(sections)
        
        # 7. 关键要点（1页）
        self.add_takeaways_slide(sections)
        
        # 8. 行动清单（1页）
        self.add_action_slide(sections)
        
        # 9. 结尾页
        self.add_conclusion_slide()
        
        return self.prs


def main():
    parser = argparse.ArgumentParser(description='将论文精读报告转换为专业PPT幻灯片')
    parser.add_argument('input_report', help='输入的Markdown报告文件')
    parser.add_argument('output_pptx', help='输出的PPT文件')
    parser.add_argument('--pdf', help='论文PDF文件路径（可选，用于提取图片）')
    
    args = parser.parse_args()
    
    input_file = args.input_report
    output_file = args.output_pptx
    pdf_file = args.pdf
    
    if not os.path.exists(input_file):
        print(f"❌ 错误：输入文件 {input_file} 不存在")
        sys.exit(1)
    
    if pdf_file and not os.path.exists(pdf_file):
        print(f"❌ 错误：PDF文件 {pdf_file} 不存在")
        sys.exit(1)
    
    try:
        # 检查依赖
        import fitz, matplotlib, PIL
        from pptx import Presentation
    except ImportError as e:
        print(f"❌ 错误：缺少依赖库 {e}")
        print("请运行: pip3 install python-pptx PyMuPDF matplotlib Pillow")
        sys.exit(1)
    
    try:
        # 读取输入文件
        print(f"📖 正在读取报告: {input_file}")
        with open(input_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # 生成PPT
        print("🎨 正在生成PPT...")
        generator = PPTGenerator()
        sections = generator.parse_markdown_report(content)
        prs = generator.generate_pptx(sections, pdf_file)
        
        # 保存PPT
        prs.save(output_file)
        
        # 统计信息
        slide_count = len(prs.slides)
        file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
        
        # 验证图片数量
        image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # 图片类型
                    image_count += 1
        
        print(f"\n✅ PPT生成完成！")
        print(f"📊 输出文件: {output_file}")
        print(f"📄 幻灯片数量: {slide_count}页")
        print(f"🖼️  图片数量: {image_count}张")
        print(f"💾 文件大小: {file_size:.2f}MB")
        print(f"🎨 主题: 白色背景，蓝白橙配色")
        
        if slide_count < 25:
            print("⚠️  警告：幻灯片数量偏少，建议检查输入报告的完整性")
        elif slide_count > 35:
            print("⚠️  提醒：幻灯片数量较多，可能需要考虑内容精简")
        
        if image_count > 0:
            print(f"✨ 成功插入 {image_count} 张图片和公式")
        
    except Exception as e:
        print(f"❌ 生成PPT时出错: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()